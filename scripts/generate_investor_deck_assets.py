#!/usr/bin/env python3
from __future__ import annotations

import json
import subprocess
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Iterable

from arabic_reshaper import reshape
from bidi.algorithm import get_display
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.colors import HexColor, white
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PDF_DIR = ROOT / "output" / "pdf"
OUTPUT_PPTX_DIR = ROOT / "output" / "pptx"
PREVIEW_DIR = ROOT / "tmp" / "pdfs"
LOGO_PATH = ROOT / "apps" / "web" / "public" / "masar-logo.png"
BACKUP_DIR = ROOT / "backups" / "prod-users-orgs-2026-03-12T20-46-54-451Z"
CERTIFICATE_PATH = Path("/Users/abdulazizalhazmi/Downloads/وثيقة العمل الحر.pdf")

PDF_PATH = OUTPUT_PDF_DIR / "masar-investor-deck-final.pdf"
PPTX_PATH = OUTPUT_PPTX_DIR / "masar-investor-deck-final.pptx"
PDF_PREVIEW_PREFIX = PREVIEW_DIR / "masar-investor-deck-final"

PDF_FONT_PATH = Path("/System/Library/Fonts/Supplemental/Arial Unicode.ttf")
PDF_FONT_NAME = "ArialUnicode"
PPTX_FONT_NAME = "Arial Unicode MS"

NAVY = "#0B1F3B"
NAVY_DARK = "#071426"
SLATE_950 = "#0F172A"
SLATE_800 = "#1E293B"
SLATE_700 = "#334155"
SLATE_600 = "#475569"
SLATE_500 = "#64748B"
SLATE_300 = "#CBD5E1"
SLATE_200 = "#E2E8F0"
SLATE_100 = "#F1F5F9"
LIGHT_BG = "#F8FAFC"
EMERALD = "#16A34A"
EMERALD_DARK = "#166534"
EMERALD_SOFT = "#DCFCE7"
BORDER = "#E2E8F0"
WHITE = "#FFFFFF"


@dataclass
class TractionSnapshot:
    total_workspaces: int
    total_users: int
    active_trials: int
    active_subscription_records: int
    trial_subscription_records: int
    owners: int
    lawyers: int
    first_workspace_date_ar: str
    snapshot_date_ar: str
    source_note: str


AR_MONTHS = {
    1: "يناير",
    2: "فبراير",
    3: "مارس",
    4: "أبريل",
    5: "مايو",
    6: "يونيو",
    7: "يوليو",
    8: "أغسطس",
    9: "سبتمبر",
    10: "أكتوبر",
    11: "نوفمبر",
    12: "ديسمبر",
}

FOUNDER_NAME = "عبدالعزيز فهد عطية الحازمي"
FOUNDER_TITLE = "المؤسس والمدير التنفيذي"
CERTIFICATE_ID = "FL-665098602"
CERTIFICATE_ROLE = "المساعدة الإدارية"
CERTIFICATE_CATEGORY = "الخدمات التخصصية"
CERTIFICATE_EXPIRY = "15 مارس 2027"
CERTIFICATE_SOURCE = "وثيقة العمل الحر صادرة من وزارة الموارد البشرية والتنمية الاجتماعية"

RAISE_AMOUNT = "2.5 مليون ريال سعودي"
INVESTMENT_HORIZON = "أفق استثماري مستهدف لمدة 5 سنوات"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", ""))


def hexcolor(hex_color: str) -> HexColor:
    return HexColor(hex_color)


def rtl(text: str) -> str:
    return get_display(reshape(text.replace("—", "-")))


def to_ar_date(date_value: datetime) -> str:
    return f"{date_value.day} {AR_MONTHS[date_value.month]} {date_value.year}"


def parse_iso(value: str) -> datetime:
    normalized = value.replace("Z", "+00:00")
    for sign in ("+", "-"):
        marker = normalized.rfind(sign)
        if marker > 10 and "." in normalized[:marker]:
            left = normalized[:marker]
            tz = normalized[marker:]
            date_part, fraction = left.split(".", 1)
            normalized = f"{date_part}.{(fraction + '000000')[:6]}{tz}"
            break
    return datetime.fromisoformat(normalized)


def load_traction() -> TractionSnapshot:
    return TractionSnapshot(
        total_workspaces=38,
        total_users=44,
        active_trials=35,
        active_subscription_records=3,
        trial_subscription_records=35,
        owners=0,
        lawyers=0,
        first_workspace_date_ar="آخر 30 يوم",
        snapshot_date_ar="مارس 2026",
        source_note="إحصائيات مباشرة من لوحة الإدارة داخل المنصة كما ظهرت في اللقطة المرفقة (آخر 30 يوم).",
    )


def register_pdf_font() -> None:
    if not PDF_FONT_PATH.exists():
        raise FileNotFoundError(f"PDF font not found: {PDF_FONT_PATH}")
    try:
        pdfmetrics.getFont(PDF_FONT_NAME)
    except KeyError:
        pdfmetrics.registerFont(TTFont(PDF_FONT_NAME, str(PDF_FONT_PATH)))


def pdf_width(text: str, size: float) -> float:
    return pdfmetrics.stringWidth(rtl(text), PDF_FONT_NAME, size)


def wrap_rtl(text: str, max_width: float, size: float) -> list[str]:
    lines: list[str] = []
    for paragraph in text.split("\n"):
        paragraph = paragraph.strip()
        if not paragraph:
            if lines and lines[-1] != "":
                lines.append("")
            continue
        words = paragraph.split()
        current: list[str] = []
        for word in words:
            candidate = " ".join(current + [word])
            if current and pdf_width(candidate, size) > max_width:
                lines.append(" ".join(current))
                current = [word]
            else:
                current.append(word)
        if current:
            lines.append(" ".join(current))
    return lines or [""]


def trim_line(text: str, max_width: float, size: float, suffix: str = " ...") -> str:
    words = text.split()
    candidate = text
    while words and pdf_width(candidate + suffix, size) > max_width:
        words.pop()
        candidate = " ".join(words)
    return (candidate + suffix).strip()


def draw_round_box(c: canvas.Canvas, x: float, y: float, w: float, h: float, *, fill: str, stroke: str, radius: float = 24, line_width: float = 1) -> None:
    c.setLineWidth(line_width)
    c.setStrokeColor(hexcolor(stroke))
    c.setFillColor(hexcolor(fill))
    c.roundRect(x, y, w, h, radius, fill=1, stroke=1)


def draw_text_block(
    c: canvas.Canvas,
    text: str,
    x_right: float,
    y_top: float,
    max_width: float,
    *,
    size: float,
    color: str,
    leading: float | None = None,
    max_lines: int | None = None,
) -> float:
    leading = leading or size * 1.4
    lines = wrap_rtl(text, max_width, size)
    if max_lines is not None and len(lines) > max_lines:
        lines = lines[:max_lines]
        lines[-1] = trim_line(lines[-1], max_width, size)

    c.setFillColor(hexcolor(color))
    c.setFont(PDF_FONT_NAME, size)
    y = y_top
    for line in lines:
        if line:
            c.drawRightString(x_right, y, rtl(line))
        y -= leading
    return y


def draw_card_title(c: canvas.Canvas, title: str, x_right: float, y_top: float, *, size: float = 16, color: str = NAVY) -> float:
    c.setFillColor(hexcolor(color))
    c.setFont(PDF_FONT_NAME, size)
    c.drawRightString(x_right, y_top, rtl(title))
    return y_top - size * 1.5


def draw_bullets(c: canvas.Canvas, items: Iterable[str], x_right: float, y_top: float, max_width: float, *, size: float = 10.5, color: str = SLATE_600, bullet_color: str = EMERALD, item_gap: float = 8) -> float:
    y = y_top
    leading = size * 1.45
    for item in items:
        lines = wrap_rtl(item, max_width - 20, size)
        c.setFillColor(hexcolor(bullet_color))
        c.circle(x_right - 5, y - 3, 2.5, fill=1, stroke=0)
        c.setFillColor(hexcolor(color))
        c.setFont(PDF_FONT_NAME, size)
        for idx, line in enumerate(lines):
            c.drawRightString(x_right - 18, y - leading * idx, rtl(line))
        y -= leading * len(lines) + item_gap
    return y


def draw_badge(c: canvas.Canvas, text: str, x_right: float, y_top: float, *, bg: str = EMERALD_SOFT, fg: str = EMERALD_DARK, size: float = 9.5) -> None:
    text_width = pdf_width(text, size)
    w = text_width + 24
    h = size + 14
    x = x_right - w
    y = y_top - h
    draw_round_box(c, x, y, w, h, fill=bg, stroke=bg, radius=16, line_width=0)
    c.setFillColor(hexcolor(fg))
    c.setFont(PDF_FONT_NAME, size)
    c.drawRightString(x_right - 12, y + 6, rtl(text))


def draw_logo_pdf(c: canvas.Canvas, x: float, y: float, w: float, h: float) -> None:
    if LOGO_PATH.exists():
        c.drawImage(ImageReader(str(LOGO_PATH)), x, y, width=w, height=h, preserveAspectRatio=True, mask="auto")


def pdf_footer(c: canvas.Canvas, slide_number: int, total: int, width: float) -> None:
    c.setStrokeColor(hexcolor(BORDER))
    c.setLineWidth(0.8)
    c.line(42, 28, width - 42, 28)
    c.setFillColor(hexcolor(SLATE_500))
    c.setFont(PDF_FONT_NAME, 8.5)
    c.drawString(42, 13, rtl("مسار المحامي | عرض استثماري"))
    c.drawRightString(width - 42, 13, rtl(f"{slide_number} / {total}"))


def pdf_metric(c: canvas.Canvas, x: float, y: float, w: float, h: float, *, value: str, label: str, note: str, dark: bool = False) -> None:
    fill = SLATE_950 if dark else WHITE
    stroke = SLATE_800 if dark else BORDER
    value_color = WHITE if dark else NAVY
    label_color = SLATE_200 if dark else NAVY
    note_color = SLATE_300 if dark else SLATE_500
    draw_round_box(c, x, y, w, h, fill=fill, stroke=stroke, radius=18)
    c.setFillColor(hexcolor(value_color))
    c.setFont(PDF_FONT_NAME, 21)
    c.drawRightString(x + w - 18, y + h - 26, rtl(value))
    draw_text_block(c, label, x + w - 18, y + h - 52, w - 36, size=10.5, color=label_color, max_lines=2)
    draw_text_block(c, note, x + w - 18, y + 26, w - 36, size=8.5, color=note_color, leading=12, max_lines=3)


def slide_cover(c: canvas.Canvas, width: float, height: float, traction: TractionSnapshot) -> None:
    c.setFillColor(hexcolor(NAVY_DARK))
    c.rect(0, 0, width, height, fill=1, stroke=0)
    c.setFillColor(hexcolor(NAVY))
    c.rect(0, height - 78, width, 78, fill=1, stroke=0)
    c.setFillColor(hexcolor(EMERALD))
    c.rect(width - 170, 0, 170, height, fill=1, stroke=0)
    draw_logo_pdf(c, 48, height - 60, 110, 34)
    draw_badge(c, "عرض استثماري", width - 204, height - 36, bg="#0E2A4F", fg="#D1FAE5")
    c.setFillColor(white)
    c.setFont(PDF_FONT_NAME, 30)
    c.drawRightString(width - 212, height - 122, rtl("مسار المحامي"))
    draw_text_block(
        c,
        "منصة عربية متخصصة لإدارة وتشغيل مكاتب المحاماة داخل السوق السعودي.",
        width - 212,
        height - 158,
        480,
        size=18.5,
        color=SLATE_100,
        leading=27,
        max_lines=3,
    )
    draw_text_block(
        c,
        "عرض مبني على ما هو موجود فعليًا داخل المنتج: إدارة القضايا والعملاء والمستندات والفوترة، ولوحة إدارة عامة على النظام، وطلبات اشتراك، وتكاملات ناجز، وشركاء النجاح.",
        width - 212,
        height - 242,
        500,
        size=11.5,
        color=SLATE_300,
        leading=18,
        max_lines=4,
    )
    cards = [
        ("النظام الأساسي", "قضايا وعملاء ومستندات ومهام وفوترة من منصة واحدة."),
        ("لوحة الإدارة", "إحصائيات مباشرة، مكاتب، مستخدمون، وطلبات اشتراك."),
        ("التوسع الحالي", "ناجز، شركاء النجاح، ونسخة شركات للباقات الأعلى قيمة."),
    ]
    x_positions = [42, 264, 486]
    for x, (title, body) in zip(x_positions, cards):
        draw_round_box(c, x, 72, 196, 116, fill="#0F1F37", stroke="#183356", radius=22)
        draw_card_title(c, title, x + 196 - 16, 154, size=13, color=EMERALD)
        draw_text_block(c, body, x + 196 - 16, 126, 160, size=9.4, color=SLATE_300, leading=14, max_lines=4)


def slide_problem(c: canvas.Canvas, width: float, height: float) -> None:
    c.setFillColor(hexcolor(LIGHT_BG))
    c.rect(0, 0, width, height, fill=1, stroke=0)
    draw_badge(c, "المشكلة والفرصة", width - 46, height - 34)
    c.setFillColor(hexcolor(NAVY))
    c.setFont(PDF_FONT_NAME, 24)
    c.drawRightString(width - 46, height - 84, rtl("لماذا يستحق هذا السوق حلًا متخصصًا؟"))
    draw_text_block(
        c,
        "القيمة لا تكمن في رقمنة أداة واحدة، بل في تحويل المكتب كاملًا إلى نظام تشغيلي منضبط.",
        width - 46,
        height - 136,
        470,
        size=11.2,
        color=SLATE_600,
        leading=18,
        max_lines=2,
    )
    draw_round_box(c, 42, 112, 400, 316, fill=WHITE, stroke=BORDER, radius=26)
    draw_card_title(c, "الأطروحة الاستثمارية", 42 + 400 - 22, 396, size=18)
    draw_bullets(
        c,
        [
            "الفئة القانونية حساسة وتحتاج ثقة وصلاحيات وعزل بيانات من البداية.",
            "الأدوات العامة لا تعكس اللغة التشغيلية المحلية ولا تكاملات السوق السعودي.",
            "المنتج الحالي تجاوز مرحلة الفكرة إلى قاعدة جاهزة لاختبار السوق والاحتفاظ.",
        ],
        42 + 400 - 22,
        350,
        350,
        size=10.7,
        item_gap=10,
    )
    cards = [
        ("الشريحة المستهدفة", "مكاتب المحاماة الصغيرة والمتوسطة أولًا، ثم نسخة الشركات الأعلى قيمة."),
        ("مبررات التوقيت", "النسخة الحالية تملك منتجًا فعليًا، وتسعيرًا جاهزًا، ووثائق إطلاق وتشغيل."),
        ("النتيجة الاستثمارية القريبة", "الانتقال من الجاهزية التقنية إلى إثبات التبني والإيراد داخل السوق."),
    ]
    y_positions = [334, 222, 110]
    for (title, body), y in zip(cards, y_positions):
        draw_round_box(c, 476, y, 438, 96, fill=WHITE, stroke=BORDER, radius=22)
        draw_card_title(c, title, 476 + 438 - 18, y + 68, size=14.5)
        draw_text_block(c, body, 476 + 438 - 18, y + 42, 384, size=10.3, color=SLATE_600, leading=15, max_lines=3)


def slide_product(c: canvas.Canvas, width: float, height: float) -> None:
    c.setFillColor(hexcolor("#F6F9FC"))
    c.rect(0, 0, width, height, fill=1, stroke=0)
    draw_badge(c, "ما بُني", width - 46, height - 34)
    c.setFillColor(hexcolor(NAVY))
    c.setFont(PDF_FONT_NAME, 24)
    c.drawRightString(width - 46, height - 84, rtl("مكونات الحل الحالية"))
    draw_text_block(
        c,
        "المحتوى أدناه يعكس الوحدات الظاهرة فعليًا داخل التطبيق ولوحة الإدارة، لا تصورات مستقبلية بعيدة.",
        width - 46,
        height - 136,
        470,
        size=11,
        color=SLATE_600,
        max_lines=3,
    )
    cards = [
        ("تشغيل المكتب", ["العملاء والقضايا والمستندات والمهام داخل مسار عمل واحد.", "الفوترة والعروض والمتابعة ضمن نفس السياق التشغيلي."]),
        ("لوحة الإدارة العامة", ["نظرة عامة على النظام مع بطاقات مباشرة للمكاتب والمستخدمين والاشتراكات.", "عرض لنسخ التجربة وعدد الفرق داخل كل مكتب."]),
        ("الإدارة المركزية", ["طلبات الاشتراك، المستخدمون، المكاتب، وفتح التبويبات الإدارية مباشرة.", "سجل تدقيق وتشغيل إداري أوضح لمدير الموقع."]),
        ("التكامل والتوسع", ["تكامل ناجز، تبويب شركاء النجاح، وبنية جاهزة للباقات الأعلى قيمة.", "ما هو ظاهر الآن يدعم التوسع التجاري بشكل تدريجي وواقعي."]),
    ]
    x_start = 42
    y_start = 260
    card_w = 420
    card_h = 126
    gap_x = 18
    gap_y = 16
    for idx, (title, bullets) in enumerate(cards):
        row = idx // 2
        col = idx % 2
        x = x_start + col * (card_w + gap_x)
        y = y_start - row * (card_h + gap_y)
        draw_round_box(c, x, y, card_w, card_h, fill=WHITE, stroke=BORDER, radius=22)
        draw_card_title(c, title, x + card_w - 18, y + card_h - 24, size=14.5)
        draw_bullets(c, bullets, x + card_w - 18, y + card_h - 56, card_w - 36, size=10.2, item_gap=8)


def slide_traction(c: canvas.Canvas, width: float, height: float, traction: TractionSnapshot) -> None:
    c.setFillColor(hexcolor(LIGHT_BG))
    c.rect(0, 0, width, height, fill=1, stroke=0)
    draw_badge(c, "نظرة عامة على النظام", width - 46, height - 34)
    c.setFillColor(hexcolor(NAVY))
    c.setFont(PDF_FONT_NAME, 24)
    c.drawRightString(width - 46, height - 86, rtl("إحصائيات مباشرة من لوحة الإدارة"))
    draw_text_block(
        c,
        "الأرقام التالية مأخوذة من لوحة الإدارة داخل المنصة، وتعكس ما يراه مدير الموقع فعليًا في النسخة الحالية خلال آخر 30 يومًا.",
        width - 46,
        height - 136,
        500,
        size=10.8,
        color=SLATE_600,
        leading=17,
        max_lines=4,
    )
    metrics = [
        (str(traction.total_workspaces), "إجمالي المكاتب", "عدد المكاتب الظاهر في لوحة الإدارة."),
        (str(traction.total_users), "إجمالي المستخدمين", "عدد الحسابات الظاهرة في لوحة الإدارة."),
        (str(traction.active_subscription_records), "الاشتراكات الفعالة", "الاشتراكات الحالية الظاهرة كخطط مفعلة."),
        (str(traction.active_trials), "النسخ التجريبية", "عدد المكاتب في مرحلة التجربة أو بانتظار التفعيل."),
    ]
    x_start = 42
    y_top = 250
    card_w = 202
    card_h = 100
    gap_x = 14
    for idx, (value, label, note) in enumerate(metrics):
        x = x_start + idx * (card_w + gap_x)
        pdf_metric(c, x, y_top, card_w, card_h, value=value, label=label, note=note)

    draw_round_box(c, 42, 124, 426, 86, fill=WHITE, stroke=BORDER, radius=20)
    draw_card_title(c, "إدارة الفرق والمكاتب", 42 + 426 - 18, 186, size=14)
    draw_text_block(c, "لوحة الإدارة تعرض أيضًا عدد فريق العمل لكل مكتب، مع حالة كل مكتب وإمكانية الانتقال مباشرة إلى تبويب المكاتب.", 42 + 426 - 18, 156, 386, size=9.8, color=SLATE_600, leading=14, max_lines=3)

    draw_round_box(c, 488, 124, 426, 86, fill=WHITE, stroke=BORDER, radius=20)
    draw_card_title(c, "تبويبات تشغيلية فعلية", 488 + 426 - 18, 186, size=14)
    draw_text_block(c, "تضم الواجهة الحالية: طلبات الاشتراك، المستخدمين، المكاتب، ناجز، شركاء النجاح، وسجل التدقيق.", 488 + 426 - 18, 156, 386, size=9.8, color=SLATE_600, leading=14, max_lines=3)

    draw_round_box(c, 42, 40, 872, 62, fill=SLATE_950, stroke=SLATE_800, radius=22)
    draw_text_block(
        c,
        "القراءة الاستثمارية: هذه مؤشرات تشغيلية حية من داخل النظام نفسه، وهي أقرب دليل على وجود منتج فعلي قيد الاستخدام، لا مجرد عرض نظري أو صفحات تسويقية فقط.",
        42 + 872 - 20,
        78,
        800,
        size=10.6,
        color=SLATE_300,
        leading=16,
        max_lines=2,
    )


def slide_model(c: canvas.Canvas, width: float, height: float) -> None:
    c.setFillColor(hexcolor("#F7FAFC"))
    c.rect(0, 0, width, height, fill=1, stroke=0)
    draw_badge(c, "الإيراد والميزة", width - 46, height - 34)
    c.setFillColor(hexcolor(NAVY))
    c.setFont(PDF_FONT_NAME, 24)
    c.drawRightString(width - 46, height - 84, rtl("نموذج الإيراد وعوامل التميز"))
    draw_text_block(
        c,
        "القيمة الاستثمارية تنبع من ارتباط التسعير بالاحتفاظ، والتوسع داخل الحساب، والمواءمة المحلية مع السوق.",
        width - 46,
        height - 136,
        470,
        size=11,
        color=SLATE_600,
        max_lines=3,
    )
    draw_round_box(c, 42, 92, 396, 320, fill=SLATE_950, stroke=SLATE_800, radius=28)
    draw_card_title(c, "التسعير الحالي", 42 + 396 - 20, 386, size=18, color=WHITE)
    pricing = [
        ("المحامي المستقل", "250 ريال", "مقعد واحد"),
        ("مكتب صغير", "500 ريال", "من 2 إلى 5 مقاعد"),
        ("مكتب متوسط", "750 ريال", "من 6 إلى 10 مقاعد"),
        ("نسخة الشركات", "تسعير تفاوضي", "من 11 إلى 30 مقعدًا"),
    ]
    px = 60
    py = 252
    pw = 168
    ph = 102
    gx = 12
    gy = 12
    for idx, (title, price, seats) in enumerate(pricing):
        row = idx // 2
        col = idx % 2
        x = px + col * (pw + gx)
        y = py - row * (ph + gy)
        draw_round_box(c, x, y, pw, ph, fill="#111E33", stroke="#203350", radius=18)
        draw_card_title(c, title, x + pw - 16, y + ph - 22, size=12.5, color=WHITE)
        c.setFillColor(hexcolor(EMERALD))
        c.setFont(PDF_FONT_NAME, 18)
        c.drawRightString(x + pw - 16, y + ph - 50, rtl(price))
        draw_text_block(c, seats, x + pw - 16, y + ph - 74, pw - 32, size=9.2, color=SLATE_300, max_lines=1)

    cards = [
        (
            "مسارات النمو",
            [
                "اشتراك مباشر منخفض الاحتكاك عبر الباقات الحالية.",
                "نمو عبر الشركاء والإحالات المرتبطة بالدفع الناجح.",
                "نسخة الشركات تفتح قيمة أعلى من الباقات الذاتية.",
            ],
        ),
        (
            "عوامل التميز",
            [
                "واجهة عربية أصلية ولغة تشغيلية ملائمة للسوق.",
                "عزل بيانات وصلاحيات وتكامل محلي يزيد صعوبة الاستبدال.",
                "ترابط التشغيل والفوترة والبوابة يرفع الاحتفاظ داخل الحساب.",
            ],
        ),
    ]
    positions = [(470, 234), (470, 78)]
    sizes = [(444, 150), (444, 130)]
    for (title, bullets), (x, y), (w, h) in zip(cards, positions, sizes):
        draw_round_box(c, x, y, w, h, fill=WHITE, stroke=BORDER, radius=22)
        draw_card_title(c, title, x + w - 18, y + h - 24, size=14.2)
        draw_bullets(c, bullets, x + w - 18, y + h - 56, w - 36, size=10.1, color=SLATE_600, item_gap=7)


def slide_founder(c: canvas.Canvas, width: float, height: float) -> None:
    c.setFillColor(hexcolor(LIGHT_BG))
    c.rect(0, 0, width, height, fill=1, stroke=0)
    draw_badge(c, "المؤسس والفريق", width - 46, height - 34)
    c.setFillColor(hexcolor(NAVY))
    c.setFont(PDF_FONT_NAME, 24)
    c.drawRightString(width - 46, height - 84, rtl("القيادة الحالية"))
    draw_text_block(
        c,
        "المشروع يقوده المؤسس حاليًا، والتمويل المقترح هدفه توسيع الفريق وبناء مسار تجاري أكثر اتساعًا.",
        width - 46,
        height - 138,
        500,
        size=10.8,
        color=SLATE_600,
        leading=17,
        max_lines=2,
    )
    draw_round_box(c, 42, 96, 414, 340, fill=SLATE_950, stroke=SLATE_800, radius=28)
    draw_card_title(c, FOUNDER_NAME, 42 + 414 - 20, 402, size=22, color=WHITE)
    draw_text_block(c, FOUNDER_TITLE, 42 + 414 - 20, 372, 330, size=12, color="#A7F3D0", max_lines=1)
    draw_bullets(
        c,
        [
            "يقود التأسيس والمنتج والتنفيذ الحالي للمشروع.",
            "يبني العرض الحالي على منتج فعلي ووثائق تشغيل واختبار وإطلاق.",
            "النموذج الحالي رشيق وعملي، وهو مناسب لمرحلة استثمارية مبكرة.",
        ],
        42 + 414 - 20,
        328,
        340,
        size=10.4,
        color=SLATE_200,
        bullet_color="#34D399",
        item_gap=10,
    )
    cards = [
        ("المرجعية المهنية", "وثيقة العمل الحر"),
        ("التخصص المسجل", CERTIFICATE_ROLE),
        ("سريان الوثيقة", f"سارية حتى {CERTIFICATE_EXPIRY}"),
        ("وضع المشروع", "يقوده المؤسس ويستعد للتوسع المدروس"),
    ]
    x_start = 484
    y_start = 240
    card_w = 208
    card_h = 108
    gap_x = 14
    gap_y = 14
    for idx, (title, body) in enumerate(cards):
        row = idx // 2
        col = idx % 2
        x = x_start + col * (card_w + gap_x)
        y = y_start - row * (card_h + gap_y)
        draw_round_box(c, x, y, card_w, card_h, fill=WHITE, stroke=BORDER, radius=20)
        draw_card_title(c, title, x + card_w - 16, y + card_h - 22, size=12.5)
        draw_text_block(c, body, x + card_w - 16, y + card_h - 50, card_w - 32, size=10, color=SLATE_600, leading=14, max_lines=3)
    draw_text_block(
        c,
        "وردت بيانات وثيقة العمل الحر هنا لأغراض التعريف المهني فقط، دون عرض أي بيانات شخصية حساسة.",
        42 + 872,
        92,
        872,
        size=9,
        color=SLATE_500,
        max_lines=2,
    )


def slide_raise(c: canvas.Canvas, width: float, height: float) -> None:
    c.setFillColor(hexcolor(NAVY_DARK))
    c.rect(0, 0, width, height, fill=1, stroke=0)
    c.setFillColor(hexcolor(NAVY))
    c.rect(0, height - 80, width, 80, fill=1, stroke=0)
    draw_badge(c, "الطرح المقترح", width - 46, height - 36, bg="#0E2A4F", fg="#D1FAE5")
    c.setFillColor(white)
    c.setFont(PDF_FONT_NAME, 26)
    c.drawRightString(width - 46, height - 94, rtl("الاستثمار المقترح وأثره"))
    c.setFillColor(hexcolor(EMERALD))
    c.setFont(PDF_FONT_NAME, 30)
    c.drawRightString(width - 46, height - 148, rtl(RAISE_AMOUNT))
    draw_text_block(c, INVESTMENT_HORIZON, width - 46, height - 176, 360, size=12, color="#A7F3D0", max_lines=1)
    draw_round_box(c, 42, 92, 446, 242, fill="#101C2F", stroke="#1F3556", radius=24)
    draw_card_title(c, "استخدام الأموال", 42 + 446 - 20, 306, size=18, color=WHITE)
    uses = [
        "40% تطوير المنتج والهندسة وسير العمل القانوني الذكي.",
        "25% اكتساب العملاء الأوائل وبناء القنوات البيعية.",
        "15% تكاملات نسخة الشركات والامتثال والتشغيل المؤسسي.",
        "10% بنية تحتية واستضافة ودعم تشغيلي.",
        "10% تصميم ومساندة تشغيلية وقانونية وإدارية.",
    ]
    draw_bullets(c, uses, 42 + 446 - 20, 266, 380, size=10.4, color=SLATE_200, bullet_color="#34D399", item_gap=10)
    draw_round_box(c, 516, 92, 398, 242, fill="#101C2F", stroke="#1F3556", radius=24)
    draw_card_title(c, "المخرجات المستهدفة من التمويل", 516 + 398 - 20, 306, size=18, color=WHITE)
    targets = [
        "تحويل المرحلة الحالية إلى قصة تبنٍ موثقة بدل الاكتفاء بالجاهزية التقنية.",
        "توسيع الفريق من قيادة فردية إلى نواة تنفيذية قابلة للتكرار.",
        "إثبات قناة اكتساب أولية مع مكاتب مدفوعة ونسخة شركات قيد التفعيل.",
        "رفع جودة البيانات والاحتفاظ والمنتج قبل أي جولة أكبر لاحقًا.",
    ]
    draw_bullets(c, targets, 516 + 398 - 20, 266, 334, size=10.3, color=SLATE_200, bullet_color="#34D399", item_gap=10)
    draw_round_box(c, 42, 38, 872, 40, fill="#0D1B30", stroke="#1F3556", radius=14)
    draw_text_block(c, "الطرح الحالي مهيأ للنقاش مع المستثمرين الأوائل ويمكن ضبط هيكله بحسب نوع الجولة.", 42 + 872 - 16, 64, 820, size=9.4, color=SLATE_300, max_lines=1)


def build_pdf(traction: TractionSnapshot) -> None:
    register_pdf_font()
    OUTPUT_PDF_DIR.mkdir(parents=True, exist_ok=True)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    width, height = 960, 540
    c = canvas.Canvas(str(PDF_PATH), pagesize=(width, height))
    c.setAuthor(FOUNDER_NAME)
    c.setTitle("Masar Investor Deck Final")
    c.setSubject("Investor deck for Masar Al-Muhami")
    slides = [
        lambda: slide_cover(c, width, height, traction),
        lambda: slide_problem(c, width, height),
        lambda: slide_product(c, width, height),
        lambda: slide_traction(c, width, height, traction),
        lambda: slide_model(c, width, height),
        lambda: slide_founder(c, width, height),
        lambda: slide_raise(c, width, height),
    ]
    total = len(slides)
    for idx, slide in enumerate(slides, start=1):
        slide()
        pdf_footer(c, idx, total, width)
        c.showPage()
    c.save()
    render_pdf_preview(PDF_PATH, PDF_PREVIEW_PREFIX)


def add_picture(slide, path: Path, left: float, top: float, width: float | None = None, height: float | None = None) -> None:
    if not path.exists():
        return
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def add_shape(slide, shape_type, left: float, top: float, width: float, height: float, *, fill: str, line: str | None = None, radius: bool = False):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    return shape


def add_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: float = 24,
    color: str = NAVY,
    bold: bool = False,
    align=PP_ALIGN.RIGHT,
    name: str = PPTX_FONT_NAME,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_bullets_pptx(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    items: Iterable[str],
    *,
    font_size: float = 18,
    color: str = SLATE_600,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = PPTX_FONT_NAME
        run.font.size = Pt(font_size)
        run.font.color.rgb = rgb(color)
        p.space_after = Pt(6)
    return box


def add_badge_pptx(slide, text: str, left: float, top: float, width: float, *, fill: str = EMERALD_SOFT, color: str = EMERALD_DARK):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, 0.38, fill=fill, line=fill)
    add_text(slide, left + 0.08, top + 0.02, width - 0.16, 0.3, text, font_size=10, color=color, bold=False)


def pptx_footer(slide, slide_number: int, total: int):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.55, 7.08, 12.2, 0.01, fill=BORDER, line=BORDER)
    add_text(slide, 0.58, 7.11, 3.6, 0.22, "مسار المحامي | عرض استثماري", font_size=8.5, color=SLATE_500, align=PP_ALIGN.LEFT)
    add_text(slide, 12.0, 7.11, 0.6, 0.22, f"{slide_number} / {total}", font_size=8.5, color=SLATE_500)


def build_pptx(traction: TractionSnapshot) -> None:
    OUTPUT_PPTX_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def cover():
        slide = prs.slides.add_slide(blank)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=NAVY_DARK, line=NAVY_DARK)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 0.78, fill=NAVY, line=NAVY)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 11.0, 0, 2.33, 7.5, fill=EMERALD, line=EMERALD)
        add_picture(slide, LOGO_PATH, 0.55, 0.18, width=1.0)
        add_badge_pptx(slide, "عرض استثماري", 9.2, 0.16, 1.6, fill="#0E2A4F", color="#D1FAE5")
        add_text(slide, 6.35, 1.35, 4.0, 0.6, "مسار المحامي", font_size=30, color=WHITE, bold=True)
        add_text(slide, 4.0, 1.95, 6.4, 1.2, "منصة عربية متخصصة لإدارة وتشغيل مكاتب المحاماة داخل السوق السعودي.", font_size=20, color=SLATE_100)
        add_text(slide, 3.75, 3.15, 6.65, 1.0, "عرض مبني على ما هو موجود فعليًا داخل المنتج: إدارة القضايا والعملاء والمستندات والفوترة، ولوحة إدارة عامة على النظام، وطلبات اشتراك، وتكاملات ناجز، وشركاء النجاح.", font_size=11.2, color=SLATE_300)
        cards = [
            (0.6, "النظام الأساسي", "قضايا وعملاء ومستندات ومهام وفوترة من منصة واحدة."),
            (3.72, "لوحة الإدارة", "إحصائيات مباشرة، مكاتب، مستخدمون، وطلبات اشتراك."),
            (6.84, "التوسع الحالي", "ناجز، شركاء النجاح، ونسخة شركات للباقات الأعلى قيمة."),
        ]
        for left, title, body in cards:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, 5.1, 2.9, 1.35, fill="#0F1F37", line="#183356")
            add_text(slide, left + 0.18, 5.28, 2.54, 0.24, title, font_size=13.5, color=EMERALD, bold=True)
            add_text(slide, left + 0.18, 5.62, 2.54, 0.55, body, font_size=9.5, color=SLATE_300)
        return slide

    def problem():
        slide = prs.slides.add_slide(blank)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=LIGHT_BG, line=LIGHT_BG)
        add_badge_pptx(slide, "المشكلة والفرصة", 10.9, 0.18, 1.75)
        add_text(slide, 7.45, 0.8, 5.0, 0.45, "لماذا يستحق هذا السوق حلًا متخصصًا؟", font_size=22, color=NAVY, bold=True)
        add_text(slide, 6.1, 1.45, 6.35, 0.6, "القيمة لا تكمن في رقمنة أداة واحدة، بل في تحويل المكتب كاملًا إلى نظام تشغيلي منضبط.", font_size=10.8, color=SLATE_600)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.58, 1.78, 5.55, 4.15, fill=WHITE, line=BORDER)
        add_text(slide, 3.15, 2.05, 2.7, 0.3, "الأطروحة الاستثمارية", font_size=17, color=NAVY, bold=True)
        add_bullets_pptx(slide, 0.95, 2.55, 4.9, 2.8, [
            "الفئة القانونية حساسة وتحتاج ثقة وصلاحيات وعزل بيانات من البداية.",
            "الأدوات العامة لا تعكس اللغة التشغيلية المحلية ولا تكاملات السوق السعودي.",
            "المنتج الحالي تجاوز مرحلة الفكرة إلى قاعدة جاهزة لاختبار السوق والاحتفاظ.",
        ], font_size=10.5)
        cards = [
            (6.6, 4.8, "الشريحة المستهدفة", "مكاتب المحاماة الصغيرة والمتوسطة أولًا، ثم نسخة الشركات الأعلى قيمة."),
            (6.6, 3.3, "مبررات التوقيت", "النسخة الحالية تملك منتجًا فعليًا، وتسعيرًا جاهزًا، ووثائق إطلاق وتشغيل."),
            (6.6, 1.8, "النتيجة الاستثمارية القريبة", "الانتقال من الجاهزية التقنية إلى إثبات التبني والإيراد داخل السوق."),
        ]
        for left, top, title, body in cards:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, 5.75, 1.15, fill=WHITE, line=BORDER)
            add_text(slide, left + 0.25, top + 0.12, 5.2, 0.24, title, font_size=13.5, color=NAVY, bold=True)
            add_text(slide, left + 0.25, top + 0.42, 5.2, 0.45, body, font_size=10.1, color=SLATE_600)
        return slide

    def product():
        slide = prs.slides.add_slide(blank)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill="#F6F9FC", line="#F6F9FC")
        add_badge_pptx(slide, "ما بُني", 11.3, 0.18, 1.3)
        add_text(slide, 8.35, 0.8, 4.0, 0.45, "مكونات الحل الحالية", font_size=22, color=NAVY, bold=True)
        add_text(slide, 5.95, 1.45, 6.45, 0.65, "المحتوى أدناه يعكس الوحدات الظاهرة فعليًا داخل التطبيق ولوحة الإدارة، لا تصورات مستقبلية بعيدة.", font_size=10.5, color=SLATE_600)
        items = [
            ("تشغيل المكتب", "العملاء والقضايا والمستندات والمهام داخل مسار عمل واحد، مع الفوترة والعروض والمتابعة."),
            ("لوحة الإدارة العامة", "نظرة عامة على النظام مع بطاقات مباشرة للمكاتب والمستخدمين والاشتراكات والنسخ التجريبية."),
            ("الإدارة المركزية", "طلبات الاشتراك، المستخدمون، المكاتب، وسجل التدقيق، مع انتقال مباشر بين التبويبات الإدارية."),
            ("التكامل والتوسع", "تكامل ناجز، تبويب شركاء النجاح، وبنية جاهزة للباقات الأعلى قيمة."),
        ]
        positions = [(0.58, 2.1), (6.95, 2.1), (0.58, 4.42), (6.95, 4.42)]
        for (title, body), (left, top) in zip(items, positions):
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, 5.8, 1.8, fill=WHITE, line=BORDER)
            add_text(slide, left + 0.22, top + 0.16, 5.2, 0.24, title, font_size=13.2, color=NAVY, bold=True)
            add_text(slide, left + 0.22, top + 0.52, 5.2, 0.7, body, font_size=9.9, color=SLATE_600)
        return slide

    def traction_slide():
        slide = prs.slides.add_slide(blank)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=LIGHT_BG, line=LIGHT_BG)
        add_badge_pptx(slide, "نظرة عامة على النظام", 10.35, 0.18, 2.3)
        add_text(slide, 7.05, 0.82, 5.3, 0.45, "إحصائيات مباشرة من لوحة الإدارة", font_size=22, color=NAVY, bold=True)
        add_text(slide, 4.95, 1.48, 7.4, 0.8, "الأرقام التالية مأخوذة من لوحة الإدارة داخل المنصة، وتعكس ما يراه مدير الموقع فعليًا في النسخة الحالية خلال آخر 30 يومًا.", font_size=10.3, color=SLATE_600)
        metrics = [
            (str(traction.total_workspaces), "إجمالي المكاتب"),
            (str(traction.total_users), "إجمالي المستخدمين"),
            (str(traction.active_subscription_records), "الاشتراكات الفعالة"),
            (str(traction.active_trials), "النسخ التجريبية"),
        ]
        positions = [(0.58, 2.15), (3.77, 2.15), (6.96, 2.15), (10.15, 2.15)]
        for (value, label), (left, top) in zip(metrics, positions):
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, 2.88, 1.48, fill=WHITE, line=BORDER)
            add_text(slide, left + 0.16, top + 0.18, 2.5, 0.3, value, font_size=21, color=NAVY, bold=True)
            add_text(slide, left + 0.16, top + 0.68, 2.5, 0.35, label, font_size=10.2, color=SLATE_600)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.58, 4.12, 6.05, 1.15, fill=WHITE, line=BORDER)
        add_text(slide, 3.82, 4.32, 2.3, 0.24, "إدارة الفرق والمكاتب", font_size=13.2, color=NAVY, bold=True)
        add_text(slide, 0.95, 4.66, 5.3, 0.44, "لوحة الإدارة تعرض أيضًا عدد فريق العمل لكل مكتب، مع حالة كل مكتب وإمكانية الانتقال مباشرة إلى تبويب المكاتب.", font_size=9.4, color=SLATE_600)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 6.75, 4.12, 5.98, 1.15, fill=WHITE, line=BORDER)
        add_text(slide, 10.15, 4.32, 2.0, 0.24, "تبويبات تشغيلية فعلية", font_size=13.2, color=NAVY, bold=True)
        add_text(slide, 7.05, 4.66, 5.2, 0.42, "تضم الواجهة الحالية: طلبات الاشتراك، المستخدمين، المكاتب، ناجز، شركاء النجاح، وسجل التدقيق.", font_size=9.4, color=SLATE_600)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.58, 5.62, 12.15, 0.9, fill=SLATE_950, line=SLATE_800)
        add_text(slide, 0.95, 5.92, 11.35, 0.3, "القراءة الاستثمارية: هذه مؤشرات تشغيلية حية من داخل النظام نفسه، وهي أقرب دليل على وجود منتج فعلي قيد الاستخدام، لا مجرد عرض نظري فقط.", font_size=10, color=SLATE_200)
        add_text(slide, 0.6, 6.78, 12.1, 0.2, traction.source_note, font_size=8.6, color=SLATE_500)
        return slide

    def model_slide():
        slide = prs.slides.add_slide(blank)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill="#F7FAFC", line="#F7FAFC")
        add_badge_pptx(slide, "الإيراد والميزة", 10.8, 0.18, 1.85)
        add_text(slide, 7.4, 0.8, 4.9, 0.45, "نموذج الإيراد وعوامل التميز", font_size=22, color=NAVY, bold=True)
        add_text(slide, 6.1, 1.45, 6.2, 0.75, "القيمة الاستثمارية تنبع من ارتباط التسعير بالاحتفاظ، والتوسع داخل الحساب، والمواءمة المحلية مع السوق.", font_size=10.5, color=SLATE_600)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.58, 1.85, 5.0, 4.0, fill=SLATE_950, line=SLATE_800)
        add_text(slide, 3.15, 2.08, 2.0, 0.35, "التسعير الحالي", font_size=17, color=WHITE, bold=True)
        pricing = [
            ("المحامي المستقل", "250 ريال", "مقعد واحد"),
            ("مكتب صغير", "500 ريال", "من 2 إلى 5 مقاعد"),
            ("مكتب متوسط", "750 ريال", "من 6 إلى 10 مقاعد"),
            ("نسخة الشركات", "تسعير تفاوضي", "من 11 إلى 30 مقعدًا"),
        ]
        positions = [(0.82, 2.48), (3.22, 2.48), (0.82, 4.08), (3.22, 4.08)]
        for (title, price, seats), (left, top) in zip(pricing, positions):
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, 2.1, 1.26, fill="#111E33", line="#203350")
            add_text(slide, left + 0.12, top + 0.12, 1.85, 0.22, title, font_size=11.5, color=WHITE, bold=True)
            add_text(slide, left + 0.12, top + 0.43, 1.85, 0.25, price, font_size=17, color=EMERALD, bold=True)
            add_text(slide, left + 0.12, top + 0.8, 1.85, 0.22, seats, font_size=8.8, color=SLATE_300)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 6.05, 2.05, 6.35, 1.7, fill=WHITE, line=BORDER)
        add_text(slide, 10.2, 2.25, 1.8, 0.22, "مسارات النمو", font_size=13.5, color=NAVY, bold=True)
        add_bullets_pptx(slide, 6.35, 2.62, 5.6, 0.9, [
            "اشتراك مباشر منخفض الاحتكاك عبر الباقات الحالية.",
            "نمو عبر الشركاء والإحالات المرتبطة بالدفع الناجح.",
            "نسخة الشركات تفتح قيمة أعلى من الباقات الذاتية.",
        ], font_size=9.8, color=SLATE_600)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 6.05, 4.0, 6.35, 1.85, fill=WHITE, line=BORDER)
        add_text(slide, 10.2, 4.2, 1.8, 0.22, "عوامل التميز", font_size=13.5, color=NAVY, bold=True)
        add_bullets_pptx(slide, 6.35, 4.57, 5.6, 1.05, [
            "واجهة عربية أصلية ولغة تشغيلية ملائمة للسوق.",
            "عزل بيانات وصلاحيات وتكامل محلي يزيد صعوبة الاستبدال.",
            "ترابط التشغيل والفوترة والبوابة يرفع الاحتفاظ داخل الحساب.",
        ], font_size=9.8, color=SLATE_600)
        return slide

    def founder():
        slide = prs.slides.add_slide(blank)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=LIGHT_BG, line=LIGHT_BG)
        add_badge_pptx(slide, "المؤسس والفريق", 10.95, 0.18, 1.7)
        add_text(slide, 8.6, 0.8, 3.7, 0.45, "القيادة الحالية", font_size=22, color=NAVY, bold=True)
        add_text(slide, 5.85, 1.5, 6.45, 0.55, "المشروع يقوده المؤسس حاليًا، والتمويل المقترح هدفه توسيع الفريق وبناء مسار تجاري أكثر اتساعًا.", font_size=10.2, color=SLATE_600)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.58, 1.85, 5.25, 4.1, fill=SLATE_950, line=SLATE_800)
        add_text(slide, 2.4, 2.15, 3.0, 0.4, FOUNDER_NAME, font_size=20, color=WHITE, bold=True)
        add_text(slide, 4.2, 2.58, 1.2, 0.22, FOUNDER_TITLE, font_size=12, color="#A7F3D0")
        add_bullets_pptx(slide, 0.9, 3.0, 4.5, 2.0, [
            "يقود التأسيس والمنتج والتنفيذ الحالي للمشروع.",
            "يبني العرض الحالي على منتج فعلي ووثائق تشغيل واختبار وإطلاق.",
            "النموذج الحالي رشيق وعملي ومناسب لمرحلة استثمارية مبكرة.",
        ], font_size=10.4, color=SLATE_200)
        details = [
            ("المرجعية المهنية", "وثيقة العمل الحر"),
            ("التخصص المسجل", CERTIFICATE_ROLE),
            ("سريان الوثيقة", f"سارية حتى {CERTIFICATE_EXPIRY}"),
            ("وضع المشروع", "يقوده المؤسس ويستعد للتوسع المدروس"),
        ]
        positions = [(6.2, 2.35), (9.35, 2.35), (6.2, 4.02), (9.35, 4.02)]
        for (title, body), (left, top) in zip(details, positions):
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, 2.78, 1.35, fill=WHITE, line=BORDER)
            add_text(slide, left + 0.18, top + 0.14, 2.35, 0.22, title, font_size=12.2, color=NAVY, bold=True)
            add_text(slide, left + 0.18, top + 0.5, 2.35, 0.48, body, font_size=9.5, color=SLATE_600)
        add_text(slide, 0.75, 6.65, 11.8, 0.22, "وردت بيانات وثيقة العمل الحر هنا لأغراض التعريف المهني فقط، دون عرض أي بيانات شخصية حساسة.", font_size=8.8, color=SLATE_500)
        return slide

    def raise_slide():
        slide = prs.slides.add_slide(blank)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=NAVY_DARK, line=NAVY_DARK)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 0.78, fill=NAVY, line=NAVY)
        add_badge_pptx(slide, "الطرح المقترح", 11.0, 0.18, 1.55, fill="#0E2A4F", color="#D1FAE5")
        add_text(slide, 7.65, 0.82, 4.7, 0.45, "الاستثمار المقترح وأثره", font_size=22, color=WHITE, bold=True)
        add_text(slide, 8.85, 1.72, 3.45, 0.4, RAISE_AMOUNT, font_size=26, color=EMERALD, bold=True)
        add_text(slide, 8.05, 2.15, 4.25, 0.22, INVESTMENT_HORIZON, font_size=11.5, color="#A7F3D0")
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.58, 3.05, 5.65, 3.0, fill="#101C2F", line="#1F3556")
        add_text(slide, 3.6, 3.32, 2.2, 0.28, "استخدام الأموال", font_size=17, color=WHITE, bold=True)
        add_bullets_pptx(slide, 0.95, 3.8, 4.9, 2.0, [
            "40% تطوير المنتج والهندسة وسير العمل القانوني الذكي.",
            "25% اكتساب العملاء الأوائل وبناء القنوات البيعية.",
            "15% تكاملات نسخة الشركات والامتثال والتشغيل المؤسسي.",
            "10% بنية تحتية واستضافة ودعم تشغيلي.",
            "10% تصميم ومساندة تشغيلية وقانونية وإدارية.",
        ], font_size=10.1, color=SLATE_200)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 6.55, 3.05, 6.2, 3.0, fill="#101C2F", line="#1F3556")
        add_text(slide, 9.4, 3.32, 3.0, 0.28, "المخرجات المستهدفة من التمويل", font_size=17, color=WHITE, bold=True)
        add_bullets_pptx(slide, 6.9, 3.8, 5.5, 2.0, [
            "تحويل المرحلة الحالية إلى قصة تبنٍ موثقة بدل الاكتفاء بالجاهزية التقنية.",
            "توسيع الفريق من قيادة فردية إلى نواة تنفيذية قابلة للتكرار.",
            "إثبات قناة اكتساب أولية مع مكاتب مدفوعة ونسخة شركات قيد التفعيل.",
            "رفع جودة البيانات والاحتفاظ والمنتج قبل أي جولة أكبر لاحقًا.",
        ], font_size=10.1, color=SLATE_200)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.58, 6.45, 12.15, 0.48, fill="#0D1B30", line="#1F3556")
        add_text(slide, 1.0, 6.59, 11.3, 0.18, "الطرح الحالي مهيأ للنقاش مع المستثمرين الأوائل ويمكن ضبط هيكله بحسب نوع الجولة.", font_size=8.7, color=SLATE_300)
        return slide

    slide_builders = [cover, problem, product, traction_slide, model_slide, founder, raise_slide]
    total = len(slide_builders)
    for idx, builder in enumerate(slide_builders, start=1):
        slide = builder()
        pptx_footer(slide, idx, total)

    prs.save(str(PPTX_PATH))
    reopened = Presentation(str(PPTX_PATH))
    if len(reopened.slides) != total:
        raise RuntimeError("PPTX validation failed")


def render_pdf_preview(pdf_path: Path, prefix: Path) -> None:
    for png in prefix.parent.glob(f"{prefix.name}-*.png"):
        png.unlink()
    subprocess.run(["pdftoppm", "-png", str(pdf_path), str(prefix)], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)


def main() -> None:
    traction = load_traction()
    build_pdf(traction)
    build_pptx(traction)
    print(f"Created PDF: {PDF_PATH}")
    print(f"Created PPTX: {PPTX_PATH}")
    print(f"Certificate available: {CERTIFICATE_PATH.exists()}")


if __name__ == "__main__":
    main()
